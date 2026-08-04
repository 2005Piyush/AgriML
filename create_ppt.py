"""
AgriML Fertilizer Recommendation System — PowerPoint Generator
Generates a professional 20-slide deck covering architecture, ML models, data pipeline,
risk analysis, sustainability, and more.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import os

# ── Color Palette ────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD   = RGBColor(0x1A, 0x25, 0x3C)
ACCENT_1  = RGBColor(0x00, 0xC9, 0x7B)   # green
ACCENT_2  = RGBColor(0x38, 0xBD, 0xF8)   # blue
ACCENT_3  = RGBColor(0xF5, 0xA6, 0x23)   # orange
ACCENT_4  = RGBColor(0xE8, 0x4D, 0x8A)   # pink
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCB, 0xD5, 0xE1)
DIM       = RGBColor(0x8B, 0x9C, 0xB2)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helpers ──────────────────────────────────────────────────────
def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_frame(slide, left, top, width, height, items, font_size=16, color=LIGHT, icon="•"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{icon}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return tf

def slide_title(slide, title, subtitle=None):
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 title, font_size=36, color=WHITE, bold=True)
    # accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_1
    bar.line.fill.background()
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11), Inches(0.5),
                     subtitle, font_size=18, color=DIM)

def add_image(slide, img_path, left, top, width=None, height=None):
    """Embed an image if it exists; returns True on success."""
    if not os.path.exists(img_path):
        print(f"[WARN] Image not found: {img_path}")
        return False
    if width and height:
        slide.shapes.add_picture(img_path, left, top, width=width, height=height)
    elif width:
        slide.shapes.add_picture(img_path, left, top, width=width)
    elif height:
        slide.shapes.add_picture(img_path, left, top, height=height)
    else:
        slide.shapes.add_picture(img_path, left, top)
    return True

def add_connector_arrow(slide, x1, y1, x2, y2, color=None):
    """Draw a simple arrow connector between two points."""
    from pptx.util import Pt
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.width = Pt(1.5)
    connector.line.color.rgb = color if color else RGBColor(0xFF, 0xFF, 0xFF)
    return connector

def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.name = "Calibri"
                paragraph.font.color.rgb = WHITE if r > 0 else BG_DARK
                paragraph.alignment = PP_ALIGN.CENTER
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ACCENT_1
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_CARD if r % 2 == 1 else RGBColor(0x14, 0x1F, 0x33)
    return table

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(sl)
# decorative circle
c = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1), Inches(5), Inches(5))
c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x0A, 0x3D, 0x2E); c.line.fill.background()
c2 = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(5), Inches(4), Inches(4))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x0E, 0x2A, 0x47); c2.line.fill.background()

add_text_box(sl, Inches(1), Inches(1.8), Inches(10), Inches(1.2),
             "🌾  AgriML", font_size=54, color=ACCENT_1, bold=True)
add_text_box(sl, Inches(1), Inches(3.0), Inches(10), Inches(0.8),
             "AI-Powered Fertilizer Recommendation System", font_size=30, color=WHITE, bold=True)
add_text_box(sl, Inches(1), Inches(3.9), Inches(10), Inches(0.6),
             "Machine Learning  •  Climate-Aware  •  Sustainable Farming",
             font_size=18, color=DIM)
add_text_box(sl, Inches(1), Inches(5.5), Inches(10), Inches(0.5),
             "Client-Side • Browser-Based • Zero Server Dependencies",
             font_size=16, color=ACCENT_2)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
slide_title(sl, "The Problem", "Why do farmers need intelligent fertilizer recommendations?")

problems = [
    "Over-fertilization wastes money and pollutes groundwater",
    "Under-fertilization leads to poor crop yields and food insecurity",
    "Nutrient imbalance locks out uptake — too much of one nutrient blocks others",
    "Climate variability (rain, heat) changes optimal application timing & quantity",
    "Manual soil testing is expensive and infrequent for small-hold farmers",
    "No easy access to agronomist expertise in rural India"
]
add_bullet_frame(sl, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5), problems, font_size=16, icon="⚠️")

# stats card
card = add_shape_rect(sl, Inches(7.2), Inches(2.0), Inches(5.3), Inches(4.5), fill_color=BG_CARD, border_color=ACCENT_3)
stats = [
    ("30%", "of global fertilizer is wasted"),
    ("₹8,000 Cr", "annual loss to Indian farmers"),
    ("10 Crops", "supported by AgriML"),
    ("7 Fertilizer", "types covered in the model"),
]
for i, (val, label) in enumerate(stats):
    y = Inches(2.2) + Inches(i * 1.1)
    add_text_box(sl, Inches(7.6), y, Inches(2), Inches(0.5), val, font_size=30, color=ACCENT_1, bold=True)
    add_text_box(sl, Inches(9.6), y + Pt(8), Inches(2.5), Inches(0.4), label, font_size=14, color=LIGHT)

# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — Solution Overview
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
slide_title(sl, "Our Solution: AgriML", "An end-to-end ML pipeline running entirely in the browser")

boxes = [
    ("🧪", "Data Input", "Soil NPK, pH, moisture,\ntemperature, rainfall,\nhumidity, crop type", ACCENT_1),
    ("⚙️", "ML Engine", "Random Forest Classifier\n+ Gradient Boosting\nRegressor", ACCENT_2),
    ("📊", "Optimization", "Quantity, cost, and\nenvironmental impact\nscoring", ACCENT_3),
    ("🌦️", "Climate Aware", "7-day forecast adjusts\ntiming and quantity\nrecommendations", ACCENT_4),
]
box_centers_x = []
for i, (icon, title, desc, accent) in enumerate(boxes):
    x = Inches(0.6 + i * 3.1)
    box_centers_x.append(0.6 + i * 3.1 + 1.4)  # center x in inch units
    card = add_shape_rect(sl, x, Inches(2.2), Inches(2.8), Inches(4.2), border_color=accent)
    add_text_box(sl, x + Inches(0.3), Inches(2.5), Inches(2.2), Inches(0.6),
                 icon, font_size=40, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.3), Inches(3.2), Inches(2.2), Inches(0.5),
                 title, font_size=20, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.2), Inches(3.9), Inches(2.4), Inches(2.0),
                 desc, font_size=14, color=LIGHT, alignment=PP_ALIGN.CENTER)

# ── Arrow connectors between the 4 solution boxes ──
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
MID_Y = 4.3  # vertical midpoint of boxes in inches
for i in range(len(box_centers_x) - 1):
    # Arrow from right-edge of box[i] to left-edge of box[i+1]
    x_start = box_centers_x[i] + 1.4   # right edge
    x_end   = box_centers_x[i + 1] - 1.4  # left edge
    conn = sl.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x_start), Inches(MID_Y),
        Inches(x_end),   Inches(MID_Y)
    )
    conn.line.color.rgb = WHITE
    conn.line.width = Pt(2)

# Flow label at bottom
add_text_box(sl, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5),
             "→  Input Soil Data   ⟶   ML Models Predict   ⟶   Optimize Quantity & Cost   ⟶   Climate-Adjusted Output",
             font_size=14, color=DIM, alignment=PP_ALIGN.CENTER)
print("[OK] Slide 3: Solution overview + flow arrows")

# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — System Architecture Diagram
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
slide_title(sl, "System Architecture", "12 interconnected modules — 100% client-side JavaScript")

# ── Try embedding the AI-generated architecture diagram on the right half ──
arch_img = os.path.join(os.path.dirname(__file__), "system_architecture_diagram.png")
if add_image(sl, arch_img, Inches(6.5), Inches(1.7), width=Inches(6.6)):
    # Left half: module list as compact bullets
    layer_groups = [
        ("Data Layer",     ACCENT_1, ["dataset.js — 2,000 synthetic training samples",
                                       "preprocessing.js — imputation + normalization",
                                       "fertilizerModel.js — Random Forest (15 trees)",
                                       "yieldModel.js — Gradient Boosting (50 stumps)"]),
        ("Analysis Layer", ACCENT_3, ["climateEngine.js — 7-day weather forecast",
                                       "optimizer.js — qty, cost & env scoring",
                                       "riskPredictor.js — 7 risk categories",
                                       "sustainability.js — 6-component index"]),
        ("Learning Layer", RGBColor(0xA7, 0x8B, 0xFA),
                                      ["adaptiveLearning.js — correction factors",
                                       "feedbackLoop.js — farmer feedback",
                                       "soilHealth.js — degradation tracker"]),
    ]
    y_offset = 1.9
    for layer_name, accent, items in layer_groups:
        add_text_box(sl, Inches(0.5), Inches(y_offset), Inches(5.7), Inches(0.38),
                     layer_name, font_size=14, color=accent, bold=True)
        y_offset += 0.42
        add_bullet_frame(sl, Inches(0.5), Inches(y_offset), Inches(5.7), Inches(len(items) * 0.32 + 0.1),
                         items, font_size=11, color=LIGHT, icon="▸")
        y_offset += len(items) * 0.32 + 0.28
else:
    # Fallback: original module boxes layout
    modules = [
        (0.5, 2.2, "dataset.js\n(Data Gen)", ACCENT_1),
        (3.5, 2.2, "preprocessing.js\n(Imputation)", ACCENT_1),
        (6.5, 2.2, "fertilizerModel.js\n(RF Classifier)", ACCENT_2),
        (9.5, 2.2, "yieldModel.js\n(GB Regressor)", ACCENT_2),
        (0.5, 4.4, "climateEngine.js", ACCENT_3),
        (3.5, 4.4, "optimizer.js", ACCENT_3),
        (6.5, 4.4, "riskPredictor.js", ACCENT_4),
        (9.5, 4.4, "sustainability.js", ACCENT_4),
        (2.0, 6.2, "adaptiveLearning.js", RGBColor(0xA7, 0x8B, 0xFA)),
        (5.0, 6.2, "feedbackLoop.js", RGBColor(0xA7, 0x8B, 0xFA)),
        (8.0, 6.2, "soilHealth.js", RGBColor(0xA7, 0x8B, 0xFA)),
    ]
    for (x, y, label, col) in modules:
        add_shape_rect(sl, Inches(x), Inches(y), Inches(2.7), Inches(1.6), border_color=col)
        add_text_box(sl, Inches(x+0.15), Inches(y+0.15), Inches(2.4), Inches(1.3),
                     label, font_size=11, color=LIGHT, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, Inches(12.3), Inches(2.5), Inches(1), Inches(0.5), "Data\nLayer",     font_size=11, color=ACCENT_1, bold=True)
    add_text_box(sl, Inches(12.3), Inches(4.7), Inches(1), Inches(0.5), "Analysis\nLayer", font_size=11, color=ACCENT_3, bold=True)
    add_text_box(sl, Inches(12.3), Inches(6.5), Inches(1), Inches(0.5), "Learning\nLayer", font_size=11, color=RGBColor(0xA7, 0x8B, 0xFA), bold=True)
print("[OK] Slide 4: System architecture + diagram image")

# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — Model Architecture Flowchart (Image)
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
slide_title(sl, "Model Architecture Flowchart", "Complete 7-stage ML pipeline from data input to recommendation output")

# Embed the flowchart image
flowchart_path = os.path.join(os.path.dirname(__file__), "AgriML_Model_Flowchart.png")
if os.path.exists(flowchart_path):
    # Center the image on the slide with padding
    img_left = Inches(2.8)
    img_top = Inches(1.6)
    img_width = Inches(7.5)
    sl.shapes.add_picture(flowchart_path, img_left, img_top, width=img_width)
    print("[OK] Flowchart image embedded in slide 5")
else:
    # Fallback: add a placeholder text box
    add_text_box(sl, Inches(2), Inches(3), Inches(9), Inches(1),
                 "⚠️  Flowchart image not found. Place 'AgriML_Model_Flowchart.png' in project root.",
                 font_size=18, color=ACCENT_3, alignment=PP_ALIGN.CENTER)
    print("[WARN] Flowchart image not found - placeholder added")

# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — Data Pipeline
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
slide_title(sl, "Data Pipeline & Preprocessing", "From raw sensor data to model-ready features")

# ── Top section: AI-generated pipeline flow image ──
pipeline_img = os.path.join(os.path.dirname(__file__), "data_pipeline_flow.png")
add_image(sl, pipeline_img, Inches(0.4), Inches(1.65), width=Inches(12.5))

# ── Bottom section: compact step cards ──
steps = [
    ("1️⃣", "Synthetic Data", "2,000 samples\n10 crops × 7 fertilizers", ACCENT_1),
    ("2️⃣", "Imputation",      "Mean/Mode\nfill missing values", ACCENT_2),
    ("3️⃣", "Normalization",   "Min-Max scaling\nto [0, 1] range", ACCENT_3),
    ("4️⃣", "One-Hot Encode",  "10 binary features\nfor crop type", ACCENT_4),
    ("5️⃣", "Feature Vector",  "18-D combined\nfeature array → model", RGBColor(0xA7, 0x8B, 0xFA)),
]
for i, (num, title, desc, accent) in enumerate(steps):
    x = Inches(0.35 + i * 2.55)
    card = add_shape_rect(sl, x, Inches(4.85), Inches(2.35), Inches(2.55), border_color=accent)
    add_text_box(sl, x + Inches(0.1), Inches(4.98), Inches(2.15), Inches(0.45),
                 num + "  " + title, font_size=13, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.1), Inches(5.5), Inches(2.15), Inches(1.6),
                 desc, font_size=12, color=LIGHT, alignment=PP_ALIGN.CENTER)
    # arrow between steps
    if i < len(steps) - 1:
        arr = sl.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            x + Inches(2.35), Inches(6.12),
            x + Inches(2.55), Inches(6.12)
        )
        arr.line.color.rgb = WHITE
        arr.line.width = Pt(1.5)
print("[OK] Slide 6: Data pipeline image + step cards")

# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — Input Features Table
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
slide_title(sl, "Input Features", "8 numeric features + 1 categorical feature (crop type)")

data = [
    ["Feature", "Type", "Unit", "Range", "Description"],
    ["Nitrogen (N)", "Numeric", "kg/ha", "0–200", "Available nitrogen in soil"],
    ["Phosphorus (P)", "Numeric", "kg/ha", "0–150", "Available phosphorus in soil"],
    ["Potassium (K)", "Numeric", "kg/ha", "0–120", "Available potassium in soil"],
    ["pH", "Numeric", "—", "3.0–10.0", "Soil acidity / alkalinity"],
    ["Moisture", "Numeric", "%", "5–100", "Soil moisture percentage"],
    ["Temperature", "Numeric", "°C", "-5–50", "Ambient temperature"],
    ["Rainfall", "Numeric", "mm", "0–450", "Monthly rainfall"],
    ["Humidity", "Numeric", "%", "10–100", "Air humidity percent"],
    ["Crop", "Categorical", "—", "10 types", "Rice, Wheat, Maize, Cotton, …"],
]
add_table(sl, Inches(0.8), Inches(2.0), Inches(11.5), Inches(5),
          len(data), 5, data,
          col_widths=[Inches(2.2), Inches(1.5), Inches(1.2), Inches(1.8), Inches(4.8)])

# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — Crop Profiles Table
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
slide_title(sl, "Supported Crop Profiles", "Optimal ranges for each crop used during training and optimization")

data = [
    ["Crop", "N (kg/ha)", "P (kg/ha)", "K (kg/ha)", "pH", "Temp (°C)", "Yield (t/ha)"],
    ["Rice",      "60–120",  "30–60",  "30–60",  "5.5–7.0",  "22–35", "3.0–8.0"],
    ["Wheat",     "80–140",  "40–70",  "20–50",  "6.0–7.5",  "12–25", "2.5–6.5"],
    ["Maize",     "80–160",  "30–60",  "20–50",  "5.8–7.0",  "18–32", "3.0–9.0"],
    ["Cotton",    "60–120",  "20–50",  "20–40",  "6.0–7.5",  "20–35", "1.5–4.0"],
    ["Sugarcane", "100–200", "40–80",  "40–80",  "5.5–7.5",  "22–36", "50–120"],
    ["Soybean",   "20–50",   "40–70",  "20–50",  "6.0–7.0",  "20–30", "1.5–4.0"],
    ["Potato",    "80–150",  "50–90",  "60–100", "5.0–6.5",  "15–25", "15–40"],
    ["Tomato",    "60–130",  "40–80",  "50–90",  "5.5–7.0",  "18–30", "20–60"],
    ["Groundnut", "10–30",   "30–60",  "20–40",  "5.5–7.0",  "22–33", "1.0–3.5"],
    ["Barley",    "60–110",  "30–55",  "20–45",  "6.0–8.0",  "10–22", "2.0–5.5"],
]
add_table(sl, Inches(0.5), Inches(1.8), Inches(12.3), Inches(5.2),
          len(data), 7, data,
          col_widths=[Inches(1.6), Inches(1.6), Inches(1.6), Inches(1.6), Inches(1.6), Inches(1.6), Inches(2.7)])

# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — Random Forest Classifier
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
slide_title(sl, "Model 1: Random Forest Classifier", "Fertilizer type recommendation using ensemble of 15 decision trees")

# Top-left — how it works (narrower to leave room for image)
card = add_shape_rect(sl, Inches(0.4), Inches(1.9), Inches(4.8), Inches(5.2), border_color=ACCENT_2)
add_text_box(sl, Inches(0.65), Inches(2.05), Inches(4.3), Inches(0.4),
             "How It Works", font_size=18, color=ACCENT_2, bold=True)
steps_rf = [
    "Bootstrap Sampling — each tree gets a random sample (with replacement)",
    "Random Feature Subset — 60% features per split for diversity",
    "Gini Impurity — measures split quality at each node",
    "Majority Vote — final prediction aggregates all 15 trees",
    "Confidence Score — % of trees agreeing on the top class",
    "Class Probabilities — derived from aggregated leaf node counts"
]
add_bullet_frame(sl, Inches(0.65), Inches(2.6), Inches(4.4), Inches(4.2), steps_rf, font_size=13, icon="→")

# Right — AI-generated Random Forest diagram image
rf_img = os.path.join(os.path.dirname(__file__), "random_forest_diagram.png")
if not add_image(sl, rf_img, Inches(5.4), Inches(1.9), width=Inches(7.6)):
    # Fallback: hyperparameters table
    card2 = add_shape_rect(sl, Inches(5.4), Inches(1.9), Inches(7.5), Inches(5.2), border_color=ACCENT_1)
    add_text_box(sl, Inches(5.7), Inches(2.05), Inches(6.5), Inches(0.4),
                 "Hyperparameters", font_size=18, color=ACCENT_1, bold=True)
    hp_data = [
        ["Parameter", "Value"],
        ["Number of Trees", "15"],
        ["Max Depth", "12"],
        ["Min Samples/Leaf", "3"],
        ["Feature Subset Ratio", "0.6 (60%)"],
        ["Split Criterion", "Gini Impurity"],
        ["Output Classes", "7 fertilizer types"],
    ]
    add_table(sl, Inches(5.7), Inches(2.7), Inches(7.0), Inches(3.8),
              len(hp_data), 2, hp_data,
              col_widths=[Inches(4.2), Inches(2.8)])
print("[OK] Slide 8: Random Forest + diagram image")

# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — Gradient Boosting Regressor
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
slide_title(sl, "Model 2: Gradient Boosting Regressor", "Yield prediction using 50 sequential decision stumps")

# Left — algorithm
card = add_shape_rect(sl, Inches(0.4), Inches(1.9), Inches(4.8), Inches(5.2), border_color=ACCENT_3)
add_text_box(sl, Inches(0.65), Inches(2.05), Inches(4.3), Inches(0.4),
             "Algorithm Steps", font_size=18, color=ACCENT_3, bold=True)
steps_gb = [
    "Initialize prediction with mean of all target yields",
    "For each of 50 estimators:",
    "   a) Compute residuals (actual – current prediction)",
    "   b) Fit a decision stump (depth 4) on residuals",
    "   c) Update predictions += learning_rate × stump output",
    "Confidence = 1 – (RMSE / target std deviation)",
    "Output: yield (t/ha), confidence, RMSE"
]
add_bullet_frame(sl, Inches(0.65), Inches(2.6), Inches(4.4), Inches(4.2), steps_gb, font_size=13, icon="→")

# Right — AI-generated Gradient Boosting diagram image
gb_img = os.path.join(os.path.dirname(__file__), "gradient_boosting_diagram.png")
if not add_image(sl, gb_img, Inches(5.4), Inches(1.9), width=Inches(7.6)):
    # Fallback: hyperparameters table
    card2 = add_shape_rect(sl, Inches(5.4), Inches(1.9), Inches(7.5), Inches(5.2), border_color=ACCENT_1)
    add_text_box(sl, Inches(5.7), Inches(2.05), Inches(6.5), Inches(0.4),
                 "Hyperparameters", font_size=18, color=ACCENT_1, bold=True)
    hp_data2 = [
        ["Parameter", "Value"],
        ["Num Estimators", "50"],
        ["Learning Rate", "0.1"],
        ["Max Depth (Stump)", "4"],
        ["Min Samples/Split", "5"],
        ["Loss Function", "MSE (Mean Squared Error)"],
        ["Output Unit", "tons / hectare"],
    ]
    add_table(sl, Inches(5.7), Inches(2.7), Inches(7.0), Inches(3.8),
              len(hp_data2), 2, hp_data2,
              col_widths=[Inches(4.2), Inches(2.8)])
print("[OK] Slide 9: Gradient Boosting + diagram image")

# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — Fertilizer Types Table
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
slide_title(sl, "Fertilizer Types & Nutrient Content", "7 fertilizer types with nutrient composition and cost data")

data_fert = [
    ["Fertilizer", "N (%)", "P (%)", "K (%)", "Cost (₹/kg)", "Best For"],
    ["Urea", "46", "0", "0", "8", "High N deficiency"],
    ["DAP", "18", "46", "0", "28", "N + P deficiency"],
    ["NPK 10-26-26", "10", "26", "26", "22", "P + K deficiency"],
    ["NPK 20-20-20", "20", "20", "20", "25", "Balanced needs"],
    ["MOP", "0", "0", "60", "18", "High K deficiency"],
    ["SSP", "0", "16", "0", "10", "Phosphorus deficiency"],
    ["Amm. Sulphate", "21", "0", "0", "12", "Acidic soils"],
]
add_table(sl, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.5),
          len(data_fert), 6, data_fert,
          col_widths=[Inches(2.2), Inches(1), Inches(1), Inches(1), Inches(1.8), Inches(4.5)])

# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — Fertilizer Selection Logic (Decision Tree Flowchart)
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
slide_title(sl, "Fertilizer Selection Logic", "Rule-based training label assignment before Random Forest training")

# ── Decision tree nodes (x, y, w, h, text, color) ──
nodes = [
    # Root decision diamond at top-centre
    (4.0, 1.9, 5.3, 1.2, "N < 50  AND  P < 35  AND  K < 30?", ACCENT_2),
    # YES branch (left)
    (0.5, 3.6, 3.2, 0.9, "✅  NPK 20-20-20\n  (Balanced Deficiency)", ACCENT_1),
    # NO branch — second decision
    (7.0, 3.6, 4.0, 0.9, "N < 50  AND  P < 35?", ACCENT_2),
    # YES → DAP
    (5.2, 5.1, 2.6, 0.9, "✅  DAP (N + P)", ACCENT_1),
    # NO → third decision
    (8.8, 5.1, 3.5, 0.9, "Only  N < 50?", ACCENT_2),
    # YES → Urea
    (7.2, 6.5, 2.2, 0.75, "✅  Urea (N only)", ACCENT_1),
    # NO → remaining rules
    (10.0, 6.5, 3.1, 0.75, "K<30→MOP  |  P<35→SSP\npH<5.5→Am.Sul  |  →NPK10-26-26", ACCENT_3),
]
for (x, y, w, h, text, color) in nodes:
    card = add_shape_rect(sl, Inches(x), Inches(y), Inches(w), Inches(h), border_color=color)
    add_text_box(sl, Inches(x + 0.1), Inches(y + 0.08), Inches(w - 0.2), Inches(h - 0.1),
                 text, font_size=11, color=LIGHT, alignment=PP_ALIGN.CENTER)

# ── YES / NO labels ──────────────────────────────────────────────
for (lx, ly, label, col) in [
    (2.6, 3.15, "YES", ACCENT_1),
    (8.6, 3.15, "NO",  ACCENT_4),
    (6.5, 4.65, "YES", ACCENT_1),
    (9.8, 4.65, "NO",  ACCENT_4),
    (7.9, 6.05, "YES", ACCENT_1),
    (10.5, 6.05, "NO", ACCENT_4),
]:
    add_text_box(sl, Inches(lx), Inches(ly), Inches(0.7), Inches(0.35),
                 label, font_size=10, color=col, bold=True, alignment=PP_ALIGN.CENTER)

# ── Legend at bottom ─────────────────────────────────────────────
add_shape_rect(sl, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4),
               fill_color=RGBColor(0x14, 0x1F, 0x33), border_color=None)
add_text_box(sl, Inches(0.6), Inches(7.05), Inches(12.0), Inches(0.35),
             "🔵 Decision node (rule check)   ✅ Green = Leaf node (fertilizer assigned)   🟠 Orange = Multiple outcomes",
             font_size=11, color=DIM, alignment=PP_ALIGN.CENTER)
print("[OK] Slide 11: Decision tree flowchart with labels")

# ══════════════════════════════════════════════════════════════════
# SLIDE 12 — Optimization Layer
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
slide_title(sl, "Optimization Layer", "Calculating optimal quantity, cost, and environmental score")

# Three columns
cols = [
    ("💊 Quantity\nOptimization", [
        "Compute N, P, K deficits vs crop profile",
        "Quantity = max deficit / nutrient % × 100",
        "Capped at 20–500 kg/hectare",
        "Scales by user's land area",
    ], ACCENT_1),
    ("💰 Cost\nEstimation", [
        "Cost = quantity × cost-per-kg",
        "Per-hectare and total cost",
        "7 fertilizers with real market rates",
        "Supports ₹ INR currency",
    ], ACCENT_3),
    ("🌍 Environmental\nImpact Score", [
        "Over-application penalty (+15)",
        "High rainfall leaching risk (+20)",
        "Soil acidification risk (+15)",
        "High moisture runoff risk (+10)",
        "Quantity-proportional impact (+0-40)",
    ], ACCENT_4),
]
for i, (title, items, accent) in enumerate(cols):
    x = Inches(0.5 + i * 4.2)
    card = add_shape_rect(sl, x, Inches(2.0), Inches(3.9), Inches(5), border_color=accent)
    add_text_box(sl, x + Inches(0.2), Inches(2.2), Inches(3.5), Inches(0.8),
                 title, font_size=18, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    add_bullet_frame(sl, x + Inches(0.2), Inches(3.2), Inches(3.5), Inches(3.5),
                     items, font_size=13, icon="•")

# ══════════════════════════════════════════════════════════════════
# SLIDE 13 — Climate Engine
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
slide_title(sl, "Climate-Aware Decision Engine", "7-day weather forecast modifies fertilizer recommendations in real-time")

# Left column — bullet points
card = add_shape_rect(sl, Inches(0.4), Inches(1.9), Inches(4.8), Inches(5.2), border_color=ACCENT_2)
add_text_box(sl, Inches(0.65), Inches(2.05), Inches(4.3), Inches(0.4),
             "Climate Logic", font_size=18, color=ACCENT_2, bold=True)
climate_items = [
    "Generates 7-day weather forecast (temp, rain, humidity, wind)",
    "Analyzes next 5 days for application window",
    "Heavy rain (>50mm) → reduce quantity 15%, delay",
    "Moderate rain (>25mm) → reduce 8%, split doses",
    "High temp (>35°C) → apply early morning",
    "High humidity (>80%) → monitor fungal risk",
    "Best Application Day — lowest risk window scored",
]
add_bullet_frame(sl, Inches(0.65), Inches(2.6), Inches(4.4), Inches(4.2), climate_items, font_size=13, icon="🌦️")

# Right — Climate engine diagram image
climate_img = os.path.join(os.path.dirname(__file__), "climate_engine_diagram.png")
if not add_image(sl, climate_img, Inches(5.4), Inches(1.9), width=Inches(7.6)):
    # Fallback: adjustment table
    card2 = add_shape_rect(sl, Inches(5.4), Inches(1.9), Inches(7.5), Inches(5.2), border_color=ACCENT_3)
    add_text_box(sl, Inches(5.7), Inches(2.05), Inches(6.5), Inches(0.4),
                 "Impact Adjustments", font_size=18, color=ACCENT_3, bold=True)
    adj_data = [
        ["Condition", "Qty Multiplier", "Risk Level"],
        ["Rain > 50mm", "0.85 (−15%)", "High"],
        ["Rain 25–50mm", "0.92 (−8%)", "Medium"],
        ["Temp > 35°C", "×0.90 (−10%)", "Medium"],
        ["Humidity > 80%", "No change", "Low"],
        ["Dry + Moderate", "1.0 (Normal)", "Low"],
    ]
    add_table(sl, Inches(5.7), Inches(2.7), Inches(7.0), Inches(3.8),
              len(adj_data), 3, adj_data,
              col_widths=[Inches(3.0), Inches(2.4), Inches(1.6)])
print("[OK] Slide 13: Climate Engine + diagram image")

# ══════════════════════════════════════════════════════════════════
# SLIDE 14 — Risk Prediction
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
slide_title(sl, "Risk Prediction Module", "7 risk categories with severity, probability, and actionable recommendations")

# Left: risk table (narrower)
risk_data = [
    ["Risk Category", "Trigger Condition", "Severity"],
    ["Over-Fertilization", "NPK excess > 20 kg/ha", "High/Med"],
    ["Nutrient Imbalance", "N:P:K ratio gap > 0.6", "High/Med"],
    ["Low Yield", "Predicted < 60% of avg", "High/Med"],
    ["pH Stress", "pH outside optimal ±1", "High/Med"],
    ["Drought Stress", "Moisture < 25%, Rain < 30", "High/Med"],
    ["Nutrient Leaching", "Rain > 200mm, Qty > 100", "High/Med"],
    ["Temperature Stress", "Temp outside optimal ±5°C", "Medium"],
]
add_table(sl, Inches(0.4), Inches(1.9), Inches(7.0), Inches(5.0),
          len(risk_data), 3, risk_data,
          col_widths=[Inches(2.5), Inches(2.8), Inches(1.7)])

# Right: AI-generated risk radar chart
risk_img = os.path.join(os.path.dirname(__file__), "risk_prediction_radar.png")
if not add_image(sl, risk_img, Inches(7.6), Inches(1.9), width=Inches(5.5)):
    # Fallback: recommendations column
    recs = [
        ["Recommendation"],
        ["Reduce application, soil test"],
        ["Apply balanced NPK first"],
        ["Review soil, ensure irrigation"],
        ["Lime (acidic) / Sulfur (alkaline)"],
        ["Irrigate, mulch"],
        ["Slow-release / split doses"],
        ["Shade / row covers"],
    ]
    add_table(sl, Inches(7.6), Inches(1.9), Inches(5.5), Inches(5.0),
              len(recs), 1, recs, col_widths=[Inches(5.5)])
print("[OK] Slide 14: Risk prediction table + radar chart")

# ══════════════════════════════════════════════════════════════════
# SLIDE 15 — Sustainability Score
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
slide_title(sl, "Sustainability Index", "6-component weighted scoring system with grades A–F")

comps = [
    ("🌿", "Fertilizer\nEfficiency", "20%", "100 − qty/5"),
    ("⚖️", "Soil\nBalance", "15%", "100 − total deficit"),
    ("🌍", "Environmental\nScore", "20%", "100 − env impact"),
    ("💧", "Water\nImpact", "15%", "Rain + moisture\n+ qty penalties"),
    ("🌾", "Soil Health\nTrend", "15%", "From degradation\ntracker"),
    ("♻️", "Carbon\nFootprint", "15%", "100 − (qty × 0.5)/2.5"),
]
for i, (icon, name, weight, formula) in enumerate(comps):
    x = Inches(0.3 + i * 2.1)
    card = add_shape_rect(sl, x, Inches(2.2), Inches(1.95), Inches(4.5), border_color=ACCENT_1)
    add_text_box(sl, x + Inches(0.1), Inches(2.4), Inches(1.75), Inches(0.5),
                 icon, font_size=30, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.1), Inches(2.9), Inches(1.75), Inches(0.7),
                 name, font_size=14, color=ACCENT_1, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.1), Inches(3.7), Inches(1.75), Inches(0.4),
                 f"Weight: {weight}", font_size=12, color=ACCENT_3, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.1), Inches(4.2), Inches(1.75), Inches(1.5),
                 formula, font_size=11, color=LIGHT, alignment=PP_ALIGN.CENTER)

# Grade legend
add_text_box(sl, Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
             "Grades:   A (≥80) Excellent   |   B (≥65) Good   |   C (≥50) Fair   |   D (≥35) Poor   |   F (<35) Critical",
             font_size=14, color=DIM, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 16 — Soil Health Tracker
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
slide_title(sl, "Soil Health Degradation Tracker", "Linear trend projection to predict nutrient decline over time")

# Left: how-it-works bullets
card = add_shape_rect(sl, Inches(0.4), Inches(1.9), Inches(4.8), Inches(5.2), border_color=ACCENT_1)
add_text_box(sl, Inches(0.65), Inches(2.05), Inches(4.3), Inches(0.4),
             "How It Works", font_size=18, color=ACCENT_1, bold=True)
soil_items = [
    "Records soil N, P, K, pH, moisture for each analysis",
    "Builds time-series history (up to 100 entries)",
    "Linear regression on nutrient values over time",
    "Predicts values 6+ months into the future",
    "Classifies trend: Declining / Stable / Improving",
    "Health score (0–100) from nutrient proximity to optimal",
    "Generates urgency-based recommendations",
]
add_bullet_frame(sl, Inches(0.65), Inches(2.6), Inches(4.4), Inches(4.2), soil_items, font_size=13, icon="📈")

# Right: AI-generated soil health trend chart
soil_img = os.path.join(os.path.dirname(__file__), "soil_health_trend.png")
if not add_image(sl, soil_img, Inches(5.4), Inches(1.9), width=Inches(7.6)):
    # Fallback: health score table
    card2 = add_shape_rect(sl, Inches(5.4), Inches(1.9), Inches(7.5), Inches(5.2), border_color=ACCENT_2)
    add_text_box(sl, Inches(5.7), Inches(2.05), Inches(6.5), Inches(0.4),
                 "Health Score Ranges", font_size=18, color=ACCENT_2, bold=True)
    health_data = [
        ["Nutrient", "Optimal Low", "Optimal High", "Score Logic"],
        ["Nitrogen",    "40 kg/ha",  "120 kg/ha", "Distance from midpoint"],
        ["Phosphorus",  "30 kg/ha",   "70 kg/ha", "Distance from midpoint"],
        ["Potassium",   "25 kg/ha",   "60 kg/ha", "Distance from midpoint"],
        ["pH",          "5.5",          "7.5",     "Distance from midpoint"],
    ]
    add_table(sl, Inches(5.7), Inches(2.7), Inches(7.0), Inches(3.5),
              len(health_data), 4, health_data,
              col_widths=[Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.9)])
print("[OK] Slide 16: Soil health + trend chart image")

# ══════════════════════════════════════════════════════════════════
# SLIDE 17 — Adaptive Learning & Feedback
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
slide_title(sl, "Adaptive Learning & Feedback Loop", "Continuous improvement through farmer feedback")

# Left column — key bullet points
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
card = add_shape_rect(sl, Inches(0.4), Inches(1.9), Inches(4.8), Inches(5.2), border_color=PURPLE)
add_text_box(sl, Inches(0.65), Inches(2.05), Inches(4.3), Inches(0.4),
             "Adaptive Learning System", font_size=16, color=PURPLE, bold=True)
adapt_items = [
    "Records every recommendation (input + predicted output)",
    "Computes correction factor: actual / predicted yield",
    "Blends: 70% model + 30% learned correction",
    "Blend weight grows with more farmer feedback (max 30%)",
    "Tracks accuracy: within 20% error = 'correct'",
    "Builds personal crop profiles over time",
]
add_bullet_frame(sl, Inches(0.65), Inches(2.6), Inches(4.4), Inches(4.2), adapt_items, font_size=12, icon="🧠")

# Right — AI-generated adaptive learning loop diagram
loop_img = os.path.join(os.path.dirname(__file__), "adaptive_learning_loop.png")
if not add_image(sl, loop_img, Inches(5.4), Inches(1.9), width=Inches(7.6)):
    # Fallback: feedback bullets
    card2 = add_shape_rect(sl, Inches(5.4), Inches(1.9), Inches(7.5), Inches(5.2), border_color=ACCENT_1)
    add_text_box(sl, Inches(5.7), Inches(2.05), Inches(6.5), Inches(0.4),
                 "Farmer Feedback System", font_size=16, color=ACCENT_1, bold=True)
    feedback_items = [
        "Farmers submit: actual yield, fertilizer used, satisfaction (1-5)",
        "System computes accuracy = 1 − |predicted − actual| / predicted",
        "Trend analysis: first-half vs second-half accuracy",
        "Per-crop accuracy breakdown",
        "Learning Status: Awaiting → Learning (1+) → Active (5+)",
        "Data persisted in localStorage (last 50 entries)",
    ]
    add_bullet_frame(sl, Inches(5.7), Inches(2.6), Inches(7.0), Inches(4.2), feedback_items, font_size=13, icon="📝")
print("[OK] Slide 17: Adaptive Learning + loop diagram image")

# ══════════════════════════════════════════════════════════════════
# SLIDE 18 — Scenario Simulator
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
slide_title(sl, "Scenario Simulator (What-If Analysis)", "Explore impact of different conditions on yield, cost, and sustainability")

scenarios = [
    ["Scenario", "Modifications", "Use Case"],
    ["📉 Reduce Fert 30%", "N, P, K × 0.7", "Cost saving analysis"],
    ["📈 Increase Fert 30%", "N, P, K × 1.3", "Yield maximization"],
    ["🏜️ Drought", "Rain=20mm, Moisture=20%", "Drought preparedness"],
    ["🌧️ Heavy Rainfall", "Rain=350mm, Moisture=85%", "Flood risk planning"],
    ["🌡️ Heat Wave", "Temp=42°C, Humidity=30%", "Extreme heat response"],
    ["✨ Optimal Conditions", "All params at ideal values", "Best-case benchmarking"],
]
add_table(sl, Inches(0.8), Inches(2.0), Inches(11.5), Inches(3.8),
          len(scenarios), 3, scenarios,
          col_widths=[Inches(3), Inches(3.5), Inches(5)])

add_text_box(sl, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.8),
             "Each scenario computes: Δ Yield (t/ha), Δ Cost (₹), Δ Environmental Impact, and generates a natural-language recommendation comparing base vs. scenario.",
             font_size=14, color=DIM)

# ══════════════════════════════════════════════════════════════════
# SLIDE 19 — Explainability & Alerts
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
slide_title(sl, "Explainability & Smart Alerts", "Every recommendation comes with human-readable reasoning")

# Left — Explainability
card = add_shape_rect(sl, Inches(0.5), Inches(2.0), Inches(5.8), Inches(5), border_color=ACCENT_2)
add_text_box(sl, Inches(0.8), Inches(2.2), Inches(5), Inches(0.4),
             "Model Explainability", font_size=20, color=ACCENT_2, bold=True)
explain_items = [
    "Feature Importance via correlation-based analysis",
    "Top-3 influencing features shown per recommendation",
    "Nutrient-specific explanations (e.g., 'Low N → use N-rich')",
    "Environmental context notes (rainfall, temperature)",
    "Yield context: above/below average for crop type",
    "Sustainability badges earned (🏆 Efficient, 🌿 Low Impact, …)",
]
add_bullet_frame(sl, Inches(0.8), Inches(2.8), Inches(5.2), Inches(4), explain_items, font_size=14, icon="💡")

# Right — Alerts
card2 = add_shape_rect(sl, Inches(6.8), Inches(2.0), Inches(5.8), Inches(5), border_color=ACCENT_3)
add_text_box(sl, Inches(7.1), Inches(2.2), Inches(5), Inches(0.4),
             "Smart Alert System", font_size=20, color=ACCENT_3, bold=True)
alerts_data = [
    ["Level", "Trigger", "Message"],
    ["🔴 Critical", "pH < 4.5", "Apply lime before fertilizing"],
    ["🔴 Critical", "Moisture < 20%", "Irrigate before application"],
    ["🟡 Warning", "pH < 5.5", "Acidic soil reduces uptake"],
    ["🟡 Warning", "pH > 8.5", "Apply gypsum or sulfur"],
    ["🟡 Warning", "N < 20", "Yellowing leaves likely"],
    ["🟡 Warning", "Temp > 40°C", "Apply early AM / evening"],
    ["🟢 Info", "Low deficits", "Minimal fertilization needed"],
]
add_table(sl, Inches(7.1), Inches(2.8), Inches(5.2), Inches(4),
          len(alerts_data), 3, alerts_data,
          col_widths=[Inches(1.2), Inches(1.5), Inches(2.5)])

# ══════════════════════════════════════════════════════════════════
# SLIDE 20 — Technology Stack & Key Metrics
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
slide_title(sl, "Technology Stack & Key Metrics")

# Left — Tech Stack
card = add_shape_rect(sl, Inches(0.5), Inches(2.0), Inches(5.8), Inches(5), border_color=ACCENT_2)
add_text_box(sl, Inches(0.8), Inches(2.2), Inches(5), Inches(0.4),
             "Technology Stack", font_size=20, color=ACCENT_2, bold=True)
tech = [
    "Frontend: Vanilla JavaScript (ES6 Modules)",
    "Styling: CSS3 with Glassmorphism design",
    "Build: Vite for bundling and dev server",
    "ML: Custom Random Forest + Gradient Boosting (no libs)",
    "Storage: localStorage for adaptive learning + history",
    "Charts: Custom SVG / Canvas charting",
    "Deployment: Static hosting (Vercel / Netlify)",
    "Zero server dependencies — runs 100% in browser",
]
add_bullet_frame(sl, Inches(0.8), Inches(2.8), Inches(5.2), Inches(4), tech, font_size=13, icon="⚡")

# Right — Metrics
card2 = add_shape_rect(sl, Inches(6.8), Inches(2.0), Inches(5.8), Inches(5), border_color=ACCENT_1)
add_text_box(sl, Inches(7.1), Inches(2.2), Inches(5), Inches(0.4),
             "Key Metrics", font_size=20, color=ACCENT_1, bold=True)
metrics_data = [
    ["Metric", "Value"],
    ["Training Samples", "2,000 synthetic"],
    ["Feature Dimensions", "18 (8 numeric + 10 encoded)"],
    ["Supported Crops", "10"],
    ["Fertilizer Classes", "7"],
    ["RF Trees", "15"],
    ["GB Estimators", "50"],
    ["Risk Categories", "7"],
    ["Sustainability Components", "6"],
    ["Code Modules", "12"],
]
add_table(sl, Inches(7.1), Inches(2.8), Inches(5.2), Inches(4),
          len(metrics_data), 2, metrics_data,
          col_widths=[Inches(3), Inches(2.2)])

# ══════════════════════════════════════════════════════════════════
# SLIDE 21 — Future Roadmap
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
slide_title(sl, "Future Roadmap", "Planned enhancements for the next phase")

roadmap = [
    ("Phase 1", "Real Weather API\nIntegration", "Connect to IMD / OpenWeather\nfor live forecast data", ACCENT_1),
    ("Phase 2", "Satellite Imagery\n(NDVI)", "Use NDVI from Sentinel-2\nfor crop health monitoring", ACCENT_2),
    ("Phase 3", "Backend + Auth", "Node.js/Firebase backend\nfor data persistence", ACCENT_3),
    ("Phase 4", "Mobile App\n(React Native)", "Offline-capable mobile\napp for rural farmers", ACCENT_4),
]
for i, (phase, title, desc, accent) in enumerate(roadmap):
    x = Inches(0.5 + i * 3.15)
    card = add_shape_rect(sl, x, Inches(2.2), Inches(2.9), Inches(4.5), border_color=accent)
    add_text_box(sl, x + Inches(0.2), Inches(2.4), Inches(2.5), Inches(0.4),
                 phase, font_size=14, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.2), Inches(2.9), Inches(2.5), Inches(0.8),
                 title, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.2), Inches(4.0), Inches(2.5), Inches(2.0),
                 desc, font_size=14, color=LIGHT, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 22 — Thank You
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
c = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(0.5), Inches(5.5), Inches(5.5))
c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x0A, 0x3D, 0x2E); c.line.fill.background()

add_text_box(sl, Inches(1), Inches(2.5), Inches(11), Inches(1),
             "Thank You!", font_size=54, color=ACCENT_1, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
             "AgriML — Empowering Farmers with AI", font_size=24, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(4.8), Inches(11), Inches(0.6),
             "Questions & Discussion", font_size=20, color=DIM, alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(6.0), Inches(11), Inches(0.4),
             "🌾  Smart Fertilizer • 📊 Yield Prediction • 🌦️ Climate Aware • ♻️ Sustainable",
             font_size=16, color=ACCENT_2, alignment=PP_ALIGN.CENTER)

# ── Add Slide Transitions & Entrance Animations ─────────────────
from lxml import etree
from copy import deepcopy

# Namespace map for OOXML
nsmap = {
    'p':   'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a':   'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r':   'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main',
}

# Transition effects to cycle through (varied for visual interest)
transition_types = [
    # (transition_element_tag, attributes, optional_child)
    ('fade', {'thruBlk': '0'}, None),
    ('push', {'dir': 'l'}, None),
    ('wipe', {'dir': 'r'}, None),
    ('cover', {'dir': 'u'}, None),
    ('fade', {'thruBlk': '1'}, None),
    ('push', {'dir': 'r'}, None),
    ('wipe', {'dir': 'd'}, None),
    ('cover', {'dir': 'l'}, None),
    ('split', {'orient': 'horz', 'dir': 'out'}, None),
    ('wheel', {'spokes': '4'}, None),
    ('fade', {'thruBlk': '0'}, None),
]

def add_slide_transition(slide, trans_type, trans_attrs, speed='med'):
    """Add a slide transition effect via XML injection."""
    slide_elem = slide._element

    # Remove existing transition if any
    existing = slide_elem.findall('{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
    for e in existing:
        slide_elem.remove(e)

    # Build transition XML
    p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    transition = etree.SubElement(slide_elem, f'{{{p_ns}}}transition')
    transition.set('spd', speed)
    transition.set('advClick', '1')

    # Add the specific transition type
    child = etree.SubElement(transition, f'{{{p_ns}}}{trans_type}')
    for key, val in trans_attrs.items():
        child.set(key, val)

def add_entrance_animations(slide):
    """Add entrance animations (appear with fade) for shapes on the slide."""
    slide_elem = slide._element
    p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # Get all shape tree children (the actual shapes)
    spTree = slide_elem.find(f'.//{{{p_ns}}}cSld/{{{p_ns}}}spTree')
    if spTree is None:
        return

    shapes = []
    for child in spTree:
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('sp', 'graphicFrame', 'grpSp', 'pic'):
            shapes.append(child)

    if not shapes:
        return

    # Build timing XML for sequential appear animations
    timing = etree.SubElement(slide_elem, f'{{{p_ns}}}timing')
    tnLst = etree.SubElement(timing, f'{{{p_ns}}}tnLst')
    par_main = etree.SubElement(tnLst, f'{{{p_ns}}}par')

    cTn_main = etree.SubElement(par_main, f'{{{p_ns}}}cTn')
    cTn_main.set('id', '1')
    cTn_main.set('dur', 'indefinite')
    cTn_main.set('restart', 'never')
    cTn_main.set('nodeType', 'tmRoot')

    childTnLst = etree.SubElement(cTn_main, f'{{{p_ns}}}childTnLst')

    # Sequence container — all shapes appear on click or after previous
    seq = etree.SubElement(childTnLst, f'{{{p_ns}}}seq')
    seq.set('concurrent', '1')
    seq.set('nextAc', 'seek')

    cTn_seq = etree.SubElement(seq, f'{{{p_ns}}}cTn')
    cTn_seq.set('id', '2')
    cTn_seq.set('dur', 'indefinite')
    cTn_seq.set('nodeType', 'mainSeq')

    childTnLst_seq = etree.SubElement(cTn_seq, f'{{{p_ns}}}childTnLst')

    # Add animation for first few shapes (not all — to keep it clean)
    anim_shapes = shapes[:min(len(shapes), 6)]
    node_id = 3

    for i, shape in enumerate(anim_shapes):
        # Get shape ID from nvSpPr or nvGraphicFramePr
        sp_id = None
        for nvPr in shape.iter():
            tag_local = nvPr.tag.split('}')[-1] if '}' in nvPr.tag else nvPr.tag
            if tag_local == 'cNvPr':
                sp_id = nvPr.get('id')
                break
        if not sp_id:
            continue

        par = etree.SubElement(childTnLst_seq, f'{{{p_ns}}}par')

        cTn_par = etree.SubElement(par, f'{{{p_ns}}}cTn')
        cTn_par.set('id', str(node_id)); node_id += 1
        cTn_par.set('fill', 'hold')

        stCondLst = etree.SubElement(cTn_par, f'{{{p_ns}}}stCondLst')
        cond = etree.SubElement(stCondLst, f'{{{p_ns}}}cond')
        cond.set('delay', '0')

        childTnLst_par = etree.SubElement(cTn_par, f'{{{p_ns}}}childTnLst')

        par2 = etree.SubElement(childTnLst_par, f'{{{p_ns}}}par')
        cTn_par2 = etree.SubElement(par2, f'{{{p_ns}}}cTn')
        cTn_par2.set('id', str(node_id)); node_id += 1
        cTn_par2.set('presetID', '10')  # Fade entrance
        cTn_par2.set('presetClass', 'entr')
        cTn_par2.set('presetSubtype', '0')
        cTn_par2.set('fill', 'hold')
        cTn_par2.set('nodeType', 'withEffect' if i > 0 else 'clickEffect')

        stCondLst2 = etree.SubElement(cTn_par2, f'{{{p_ns}}}stCondLst')
        cond2 = etree.SubElement(stCondLst2, f'{{{p_ns}}}cond')
        cond2.set('delay', str(i * 200))

        childTnLst_par2 = etree.SubElement(cTn_par2, f'{{{p_ns}}}childTnLst')

        # Set element — fade effect
        aset = etree.SubElement(childTnLst_par2, f'{{{p_ns}}}set')
        cBhvr = etree.SubElement(aset, f'{{{p_ns}}}cBhvr')
        cTn_set = etree.SubElement(cBhvr, f'{{{p_ns}}}cTn')
        cTn_set.set('id', str(node_id)); node_id += 1
        cTn_set.set('dur', '1')
        cTn_set.set('fill', 'hold')
        stCondLst3 = etree.SubElement(cTn_set, f'{{{p_ns}}}stCondLst')
        cond3 = etree.SubElement(stCondLst3, f'{{{p_ns}}}cond')
        cond3.set('delay', '0')

        tgtEl = etree.SubElement(cBhvr, f'{{{p_ns}}}tgtEl')
        spTgt = etree.SubElement(tgtEl, f'{{{p_ns}}}spTgt')
        spTgt.set('spid', sp_id)

        attrNameLst = etree.SubElement(cBhvr, f'{{{p_ns}}}attrNameLst')
        attrName = etree.SubElement(attrNameLst, f'{{{p_ns}}}attrName')
        attrName.text = 'style.visibility'

        to_elem = etree.SubElement(aset, f'{{{p_ns}}}to')
        val = etree.SubElement(to_elem, f'{{{p_ns}}}strVal')
        val.set('val', 'visible')

        # AnimEffect — fade
        animEffect = etree.SubElement(childTnLst_par2, f'{{{p_ns}}}animEffect')
        animEffect.set('transition', 'in')
        animEffect.set('filter', 'fade')

        cBhvr2 = etree.SubElement(animEffect, f'{{{p_ns}}}cBhvr')
        cTn_anim = etree.SubElement(cBhvr2, f'{{{p_ns}}}cTn')
        cTn_anim.set('id', str(node_id)); node_id += 1
        cTn_anim.set('dur', '500')

        tgtEl2 = etree.SubElement(cBhvr2, f'{{{p_ns}}}tgtEl')
        spTgt2 = etree.SubElement(tgtEl2, f'{{{p_ns}}}spTgt')
        spTgt2.set('spid', sp_id)

    # Previous button / next button for sequence
    prevCondLst = etree.SubElement(seq, f'{{{p_ns}}}prevCondLst')
    prevCond = etree.SubElement(prevCondLst, f'{{{p_ns}}}cond')
    prevCond.set('evt', 'onPrev')
    prevCond.set('delay', '0')
    tgtEl_prev = etree.SubElement(prevCond, f'{{{p_ns}}}tgtEl')
    sldTgt_prev = etree.SubElement(tgtEl_prev, f'{{{p_ns}}}sldTgt')

    nextCondLst = etree.SubElement(seq, f'{{{p_ns}}}nextCondLst')
    nextCond = etree.SubElement(nextCondLst, f'{{{p_ns}}}cond')
    nextCond.set('evt', 'onNext')
    nextCond.set('delay', '0')
    tgtEl_next = etree.SubElement(nextCond, f'{{{p_ns}}}tgtEl')
    sldTgt_next = etree.SubElement(tgtEl_next, f'{{{p_ns}}}sldTgt')


# Apply transitions and animations to all slides
for i, slide in enumerate(prs.slides):
    t_idx = i % len(transition_types)
    trans_type, trans_attrs, _ = transition_types[t_idx]
    add_slide_transition(slide, trans_type, trans_attrs)
    add_entrance_animations(slide)

print("[OK] Added slide transitions + entrance animations to all slides")

# ── Save ─────────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(__file__), "AgriML_Fertilizer_Recommendation_Presentation.pptx")
prs.save(output_path)
print(f"[OK] Presentation saved to: {output_path}")
print(f"[OK] Total slides: {len(prs.slides)}")
